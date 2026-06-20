from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
OUTPUT_FILE = PROJECT_ROOT / "output" / "reports" / "report.pptx"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# "Midnight Executive" palette -- navy dominant, ice-blue secondary, white accent.
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x2A)
MUTED_TEXT = RGBColor(0x6B, 0x70, 0x85)
ACCENT = RGBColor(0xF9, 0x61, 0x67)

TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout


def _set_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _disable_shadow(shape):
    """
    Autoshapes created by python-pptx carry a <p:style> element that
    references theme effects (effectRef), which is what actually supplies
    the default drop shadow -- shape.shadow.inherit = False and an empty
    spPr effectLst do not override it reliably across renderers. Removing
    the <p:style> element entirely is the dependable fix.
    """
    from pptx.oxml.ns import qn

    sp = shape._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                  color=DARK_TEXT, font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    if anchor is not None:
        tf.vertical_anchor = anchor

    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font

    return box


def _add_bullets(slide, left, top, width, height, bullets, *, size=16,
                  color=DARK_TEXT, font=BODY_FONT, line_spacing=1.3, vertical_center=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if vertical_center:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.line_spacing = line_spacing
        p.space_after = Pt(14)
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font

    return box


def _add_picture_fit(slide, image_path, left, top, max_width, max_height):
    """Adds a picture scaled to fit inside the given box, centered within it."""

    pic = slide.shapes.add_picture(image_path, left, top)

    width_ratio = max_width / pic.width
    height_ratio = max_height / pic.height
    scale = min(width_ratio, height_ratio, 1.0)

    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(left + (max_width - pic.width) / 2)
    pic.top = int(top + (max_height - pic.height) / 2)

    return pic


def _build_title_slide(prs, slide_data):
    slide = _blank_slide(prs)
    _set_background(slide, NAVY)

    _add_textbox(
        slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.5),
        slide_data.get("title") or "Data Analysis Report",
        size=44, bold=True, color=WHITE, font=TITLE_FONT,
    )

    subtitle = slide_data.get("subtitle")
    if subtitle:
        _add_textbox(
            slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
            subtitle, size=18, color=ICE_BLUE, font=BODY_FONT,
        )

    return slide


def _build_closing_slide(prs, slide_data):
    slide = _blank_slide(prs)
    _set_background(slide, NAVY)

    _add_textbox(
        slide, Inches(1), Inches(3.1), Inches(11.3), Inches(1.2),
        slide_data.get("title") or "Thank You",
        size=40, bold=True, color=WHITE, font=TITLE_FONT,
    )

    subtitle = slide_data.get("subtitle")
    if subtitle:
        _add_textbox(
            slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
            subtitle, size=16, color=ICE_BLUE, font=BODY_FONT,
        )

    return slide


def _build_bullets_slide(prs, slide_data):
    slide = _blank_slide(prs)
    _set_background(slide, WHITE)

    _add_textbox(
        slide, Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.9),
        slide_data.get("title") or "",
        size=32, bold=True, color=NAVY, font=TITLE_FONT,
    )

    bullets = slide_data.get("bullets") or []
    if bullets:
        # Fewer bullets get bigger type and more line spacing so a short
        # list doesn't look like an unfinished slide -- a 3-bullet slide
        # and an 8-bullet slide shouldn't use the same density.
        if len(bullets) <= 4:
            size, line_spacing = 22, 1.6
        elif len(bullets) <= 6:
            size, line_spacing = 19, 1.4
        else:
            size, line_spacing = 16, 1.25

        _add_bullets(
            slide, Inches(0.9), Inches(1.7), Inches(11.3), Inches(5.1),
            bullets, size=size, color=DARK_TEXT, line_spacing=line_spacing,
            vertical_center=True,
        )

    return slide


def _build_chart_slide(prs, slide_data):
    slide = _blank_slide(prs)
    _set_background(slide, WHITE)

    _add_textbox(
        slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        slide_data.get("title") or "",
        size=30, bold=True, color=NAVY, font=TITLE_FONT,
    )

    chart_path = slide_data.get("chart_path")
    if chart_path and Path(chart_path).exists():
        _add_picture_fit(
            slide, chart_path,
            Inches(1.2), Inches(1.5),
            Inches(10.9), Inches(5.0),
        )

    caption = slide_data.get("caption")
    if caption:
        _add_textbox(
            slide, Inches(1.2), Inches(6.7), Inches(10.9), Inches(0.5),
            caption, size=13, color=MUTED_TEXT, align=PP_ALIGN.LEFT,
        )

    return slide


def _build_two_column_slide(prs, slide_data):
    slide = _blank_slide(prs)
    _set_background(slide, WHITE)

    _add_textbox(
        slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        slide_data.get("title") or "",
        size=30, bold=True, color=NAVY, font=TITLE_FONT,
    )

    content_top = Inches(1.7)
    content_height = Inches(5.2)

    bullets = slide_data.get("bullets") or []
    if bullets:
        size = 18 if len(bullets) <= 4 else 16
        _add_bullets(
            slide, Inches(0.7), content_top, Inches(5.5), content_height,
            bullets, size=size, color=DARK_TEXT, line_spacing=1.4,
            vertical_center=True,
        )

    chart_path = slide_data.get("chart_path")
    if chart_path and Path(chart_path).exists():
        _add_picture_fit(
            slide, chart_path,
            Inches(6.5), content_top, Inches(6.0), content_height,
        )

    return slide


def _build_stat_callout_slide(prs, slide_data):
    slide = _blank_slide(prs)
    _set_background(slide, WHITE)

    _add_textbox(
        slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        slide_data.get("title") or "",
        size=30, bold=True, color=NAVY, font=TITLE_FONT,
    )

    stats = slide_data.get("stats") or []
    if not stats:
        return slide

    n = len(stats)
    card_width = Inches(11.3) / max(n, 1) - Inches(0.3)
    gap = Inches(0.3)
    start_x = Inches(1.0)
    card_top = Inches(2.9)
    card_height = Inches(3.0)

    for i, stat in enumerate(stats):
        left = int(start_x + i * (card_width + gap))

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = ICE_BLUE
        card.line.fill.background()
        card.shadow.inherit = False
        _disable_shadow(card)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)

        p_value = tf.paragraphs[0]
        p_value.text = str(stat.get("value", ""))
        p_value.alignment = PP_ALIGN.CENTER
        p_value.runs[0].font.size = Pt(34)
        p_value.runs[0].font.bold = True
        p_value.runs[0].font.color.rgb = NAVY
        p_value.runs[0].font.name = TITLE_FONT

        p_label = tf.add_paragraph()
        p_label.text = str(stat.get("label", ""))
        p_label.alignment = PP_ALIGN.CENTER
        p_label.space_before = Pt(6)
        p_label.runs[0].font.size = Pt(14)
        p_label.runs[0].font.color.rgb = MUTED_TEXT
        p_label.runs[0].font.name = BODY_FONT

    return slide


_LAYOUT_BUILDERS = {
    "title": _build_title_slide,
    "bullets": _build_bullets_slide,
    "chart": _build_chart_slide,
    "two_column": _build_two_column_slide,
    "stat_callout": _build_stat_callout_slide,
    "closing": _build_closing_slide,
}


def create_ppt(slide_plan: list[dict], output_path: str | Path | None = None) -> str:
    output_file = Path(output_path) if output_path else OUTPUT_FILE
    output_file.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if not slide_plan:
        slide_plan = [{
            "layout": "title",
            "title": "Data Analysis Report",
            "subtitle": "No slide content was generated.",
        }]

    for slide_data in slide_plan:
        layout = slide_data.get("layout", "bullets")
        builder = _LAYOUT_BUILDERS.get(layout, _build_bullets_slide)
        builder(prs, slide_data)

    prs.save(output_file)

    return str(output_file)